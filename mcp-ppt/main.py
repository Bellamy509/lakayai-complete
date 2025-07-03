#!/usr/bin/env python3
"""
PowerPoint Creator MCP Server
Serveur MCP pour créer des présentations PowerPoint avec python-pptx
"""

import asyncio
import json
import sys
from typing import Any, Dict, List, Optional
from pathlib import Path

# MCP SDK imports
from mcp.server.models import InitializationOptions
from mcp.server import NotificationOptions, Server
from mcp.server.stdio import stdio_server
from mcp.types import Resource, Tool, TextContent, ImageContent, EmbeddedResource

# PowerPoint imports
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Initialisation du serveur MCP
server = Server("PowerPoint-Creator")

@server.list_tools()
async def handle_list_tools() -> List[Tool]:
    """Liste tous les outils disponibles pour PowerPoint"""
    return [
        Tool(
            name="create_presentation",
            description="Créer une nouvelle présentation PowerPoint vide",
            inputSchema={
                "type": "object",
                "properties": {
                    "title": {
                        "type": "string",
                        "description": "Titre de la présentation"
                    },
                    "filename": {
                        "type": "string", 
                        "description": "Nom du fichier (sans extension .pptx)"
                    }
                },
                "required": ["title", "filename"]
            }
        ),
        Tool(
            name="add_title_slide",
            description="Ajouter une diapositive de titre à la présentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {
                        "type": "string",
                        "description": "Nom du fichier de présentation"
                    },
                    "title": {
                        "type": "string",
                        "description": "Titre principal de la diapositive"
                    },
                    "subtitle": {
                        "type": "string", 
                        "description": "Sous-titre (optionnel)"
                    }
                },
                "required": ["filename", "title"]
            }
        ),
        Tool(
            name="add_content_slide",
            description="Ajouter une diapositive avec titre et contenu",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {
                        "type": "string",
                        "description": "Nom du fichier de présentation"
                    },
                    "title": {
                        "type": "string",
                        "description": "Titre de la diapositive"
                    },
                    "content": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Liste des points de contenu"
                    }
                },
                "required": ["filename", "title", "content"]
            }
        ),
        Tool(
            name="add_image_slide",
            description="Ajouter une diapositive avec une image",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {
                        "type": "string",
                        "description": "Nom du fichier de présentation"
                    },
                    "title": {
                        "type": "string",
                        "description": "Titre de la diapositive"
                    },
                    "image_path": {
                        "type": "string",
                        "description": "Chemin vers l'image à insérer"
                    },
                    "caption": {
                        "type": "string",
                        "description": "Légende de l'image (optionnel)"
                    }
                },
                "required": ["filename", "title", "image_path"]
            }
        ),
        Tool(
            name="save_presentation",
            description="Sauvegarder la présentation PowerPoint",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {
                        "type": "string",
                        "description": "Nom du fichier de présentation"
                    },
                    "output_path": {
                        "type": "string",
                        "description": "Chemin de sortie (optionnel, par défaut dossier actuel)"
                    }
                },
                "required": ["filename"]
            }
        )
    ]

# Stockage des présentations en mémoire
presentations: Dict[str, Presentation] = {}

@server.call_tool()
async def handle_call_tool(name: str, arguments: Dict[str, Any]) -> List[TextContent]:
    """Gestionnaire d'appels d'outils"""
    
    if name == "create_presentation":
        filename = arguments["filename"]
        title = arguments["title"]
        
        # Créer une nouvelle présentation
        prs = Presentation()
        presentations[filename] = prs
        
        return [TextContent(
            type="text",
            text=f"✅ Présentation '{title}' créée avec succès (fichier: {filename})"
        )]
    
    elif name == "add_title_slide":
        filename = arguments["filename"]
        title = arguments["title"]
        subtitle = arguments.get("subtitle", "")
        
        if filename not in presentations:
            return [TextContent(
                type="text", 
                text=f"❌ Erreur: Présentation '{filename}' non trouvée. Créez d'abord une présentation."
            )]
        
        prs = presentations[filename]
        
        # Layout de diapositive titre (index 0)
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Ajouter le titre
        slide.shapes.title.text = title
        
        # Ajouter le sous-titre si fourni
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
        
        return [TextContent(
            type="text",
            text=f"✅ Diapositive de titre ajoutée: '{title}'"
        )]
    
    elif name == "add_content_slide":
        filename = arguments["filename"]
        title = arguments["title"]
        content = arguments["content"]
        
        if filename not in presentations:
            return [TextContent(
                type="text",
                text=f"❌ Erreur: Présentation '{filename}' non trouvée. Créez d'abord une présentation."
            )]
        
        prs = presentations[filename]
        
        # Layout avec titre et contenu (index 1)
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        # Ajouter le titre
        slide.shapes.title.text = title
        
        # Ajouter le contenu
        if len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()
            
            for i, point in enumerate(content):
                if i == 0:
                    tf.text = point
                else:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0
        
        return [TextContent(
            type="text",
            text=f"✅ Diapositive de contenu ajoutée: '{title}' avec {len(content)} points"
        )]
    
    elif name == "add_image_slide":
        filename = arguments["filename"]
        title = arguments["title"]
        image_path = arguments["image_path"]
        caption = arguments.get("caption", "")
        
        if filename not in presentations:
            return [TextContent(
                type="text",
                text=f"❌ Erreur: Présentation '{filename}' non trouvée. Créez d'abord une présentation."
            )]
        
        # Vérifier que l'image existe
        if not Path(image_path).exists():
            return [TextContent(
                type="text",
                text=f"❌ Erreur: Image non trouvée: {image_path}"
            )]
        
        prs = presentations[filename]
        
        # Layout vide (index 6) pour plus de contrôle
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Ajouter le titre manuellement
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Inches(0.4)
        title_frame.paragraphs[0].font.bold = True
        
        # Ajouter l'image
        try:
            slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(8))
            
            result_text = f"✅ Diapositive avec image ajoutée: '{title}'"
            if caption:
                # Ajouter une légende
                caption_shape = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(1))
                caption_frame = caption_shape.text_frame
                caption_frame.text = caption
                result_text += f" (avec légende)"
            
            return [TextContent(type="text", text=result_text)]
            
        except Exception as e:
            return [TextContent(
                type="text",
                text=f"❌ Erreur lors de l'ajout de l'image: {str(e)}"
            )]
    
    elif name == "save_presentation":
        filename = arguments["filename"]
        output_path = arguments.get("output_path", ".")
        
        if filename not in presentations:
            return [TextContent(
                type="text",
                text=f"❌ Erreur: Présentation '{filename}' non trouvée. Créez d'abord une présentation."
            )]
        
        prs = presentations[filename]
        
        # Construire le chemin de sortie
        output_dir = Path(output_path)
        if not output_dir.exists():
            output_dir.mkdir(parents=True, exist_ok=True)
        
        output_file = output_dir / f"{filename}.pptx"
        
        try:
            prs.save(str(output_file))
            return [TextContent(
                type="text",
                text=f"✅ Présentation sauvegardée: {output_file}"
            )]
        except Exception as e:
            return [TextContent(
                type="text",
                text=f"❌ Erreur lors de la sauvegarde: {str(e)}"
            )]
    
    else:
        return [TextContent(
            type="text",
            text=f"❌ Outil inconnu: {name}"
        )]

async def main():
    """Fonction principale"""
    # Configuration des options du serveur
    options = InitializationOptions(
        server_name="PowerPoint-Creator",
        server_version="1.0.0",
        capabilities=server.get_capabilities(
            notification_options=NotificationOptions(),
            experimental_capabilities={}
        )
    )
    
    async with stdio_server() as (read_stream, write_stream):
        await server.run(
            read_stream,
            write_stream,
            options
        )

if __name__ == "__main__":
    asyncio.run(main()) 